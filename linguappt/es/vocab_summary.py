from linguappt.ppt import PPT
from linguappt.filereader import readCSV
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os
import math
import json

class SpanishVocabMeta:

  pos_info = {
    'noun': ('名词', 'el nombre', ['noun.', 'f.', 'm.', 'f.m.', 'propn,', 'f.pl.', 'm.pl.']),
    'adj': ('形容词', 'el adjectivo', ['adj.']),
    'verb': ('动词', 'el verbo', ['verb.', 'vr.', 'vi.', 'vt.', 'aux.']),
    'adv': ('副词', 'el adverbio', ['adv.']),
    'pron': ('代词', 'los pronombres', ['pron.']),
    'prep': ('前置词', '', ['prep.', 'adp.']),
    'other': ('其他', 'los otros', [])
  }

  tense_info = {
    'imperativo_afirmativo': "命令式-肯定",
    'imperativo_negativo': "命令式-否定",
    'indicativo-pretérito': "陈述式-过去时",
    'indicativo-presente': "陈述式-现在时",
    'subjuntivo-presente': "虚拟式-现在时",
    'indicativo-futuro': "陈述式-将来时",
    'subjuntivo-futuro': "虚拟式-将来时",
    'indicativo-imperfecto': "陈述式-未完成时",
    'subjuntivo-imperfecto': "虚拟式-未完成时",
    'indicativo-condicional': "条件式",
    'participio': "过去分词",
    'gerundio': "现在分词"
  }

class SpanishVocabPPT(PPT):

  template_dir = os.path.dirname(__file__)
  templates = {
    "classic": os.path.join(template_dir, 'templates/vocab_spanish_classic.pptx'),
    "watermark": os.path.join(template_dir, 'templates/vocab_spanish_watermark.pptx')
  }
  lang = 'es'

  def __init__(self, sourcefile, title="", genre="classic"):
    """
    word_distribution is a list of dict
    e.g,
    word_distribution = [
      {word: 'Hola', pos: 'INTJ', num: 2},
      {word: 'bueno', pos: 'ADJ', num: 20},
      {word: 'pues', pos: 'ADV', num: 4},
    ]
    """

    self.template = SpanishVocabPPT.templates[genre]

    self.keys = ['num', 'word', 'pos', 'meaning', 'dict_pos', 'from', 'extension', 'variations']
    self.content = self.assert_content(sourcefile)

    self.divide_by_pos()
    self.cal_displayed_word_distribution()
    self.cal_distribution()

    self.prs = Presentation(self.template)
    self.title = title
    
  def assert_content(self, sourcefile):
    """
      Ensure the content has keys defined in ppt class
    """
    content = readCSV(sourcefile) 
    for e in content:
      keys = e.keys()
      assert(len(keys) == len(self.keys))
      for k in keys:
        assert(k in self.keys)
    return content

  def divide_by_pos(self):
    """
      Store words in self.word_distribution as dict, whose key is dict_pos and value is a list of words
      dict_pos is defined in x2cdict package
    """
    self.word_distribution  = {}

    for e in self.content:
      key = e["dict_pos"]
      if key not in self.word_distribution:
        self.word_distribution[key] = []
      self.word_distribution[key].append(e)

  def cal_displayed_word_distribution(self):
    """
      Bridge dict_pos and pos defined in SpanishVocabMeta
      Store words in self.displayed_word_distribution as dict, which are used to presented in ppt.
      Key in self.displayed_word_distribution is pos defined SpanishVocabMeta and value is a list of words
    """
    self.displayed_word_distribution = {}
    for key in SpanishVocabMeta.pos_info.keys():
      poses = SpanishVocabMeta.pos_info[key][2] 
      self.displayed_word_distribution[key] = []
      for pos in poses:
        if pos in self.word_distribution.keys():
          self.displayed_word_distribution[key].extend(self.word_distribution[pos])


  def cal_distribution(self):
    """
      Calculate the number of words according to pos defined in SpanishVocabMeta
    """
    self.distribution = [{"pos": pos, "name": SpanishVocabMeta.pos_info[pos][1], "num": len(ws)} for pos, ws in self.displayed_word_distribution.items()]

  def create_home(self):
    layout = self.prs.slide_layouts.get_by_name("Title and subtitle")
    slide = self.prs.slides.add_slide(layout)
    holders = slide.shapes.placeholders
 
    title = holders[10] 
    title.text_frame.text = self.title
    subtitle = holders[11]
    subtitle.text_frame.text = '词汇总结'

  def create_statistic(self):
    layout = self.prs.slide_layouts.get_by_name("Word count")
    slide = self.prs.slides.add_slide(layout)
    holders = slide.shapes.placeholders

    for index, info in enumerate(self.distribution[:3]):
      pos = holders[11+2*index]
      num = holders[10+2*index]
      pos.text_frame.text = info["name"].upper() 
      num.text_frame.text = str(info["num"])


  def create_vocab_title(self, index, title_content, subtitle_content):
    layout = self.prs.slide_layouts.get_by_name('Title '+str(index))
    slide = self.prs.slides.add_slide(layout)
    holders = slide.shapes.placeholders

    title = holders[10]
    title.text_frame.text = title_content
    subtitle = holders[11]
    subtitle.text_frame.text = subtitle_content
    note = slide.notes_slide
    note.notes_text_frame.text = title_content.lower()


  def create_noun_m_word(self, v):
    layout = self.prs.slide_layouts.get_by_name("Noun m vocab")
    slide = self.prs.slides.add_slide(layout)
    holders = slide.shapes.placeholders

    noun, meaning = holders[10], holders[11]
    noun.text_frame.text = v["word"]
    ms = v["meaning"].split(",")
    if len(ms) > 4:
      ms = ms[:4]
    meaning.text_frame.text = "\n".join(ms) 

    s_def, s = holders[12], holders[13]
    s_def.text_frame.text = "el"
    s.text_frame.text = v["word"]

    s_undef, s = holders[14], holders[15]
    s_undef.text_frame.text = "un"
    s.text_frame.text = v["word"]

    if v["extension"] != "":
      extension = json.loads(v["extension"].replace("\'", "\""))

      pl_def, pl = holders[16], holders[17]
      pl_def.text_frame.text = "los"
      pl.text_frame.text = extension["mpl"]

      pl_undef, pl = holders[18], holders[19]
      pl_undef.text_frame.text = "unos"
      pl.text_frame.text = extension["mpl"]

    note = slide.notes_slide
    note.notes_text_frame.text = v["word"]

  def create_noun_f_word(self, v):
    layout = self.prs.slide_layouts.get_by_name("Noun f vocab")
    slide = self.prs.slides.add_slide(layout)
    holders = slide.shapes.placeholders

    noun, meaning = holders[10], holders[11]
    noun.text_frame.text = v["word"]
    ms = v["meaning"].split(",")
    if len(ms) > 4:
      ms = ms[:4]
    meaning.text_frame.text = "\n".join(ms) 

    s_def, s = holders[12], holders[13]
    s_def.text_frame.text = "la"
    s.text_frame.text = v["word"]

    s_undef, s = holders[14], holders[15]
    s_undef.text_frame.text = "una"
    s.text_frame.text = v["word"]

    if v["extension"] != "":
      extension = json.loads(v["extension"].replace("\'", "\""))

      pl_def, pl = holders[16], holders[17]
      pl_def.text_frame.text = "las"
      pl.text_frame.text = extension["fpl"]

      pl_undef, pl = holders[18], holders[19]
      pl_undef.text_frame.text = "unas"
      pl.text_frame.text = extension["fpl"]

    note = slide.notes_slide
    note.notes_text_frame.text = v["word"]

  def create_noun_mpl_word(self, v):
    layout = self.prs.slide_layouts.get_by_name("Noun m vocab")
    slide = self.prs.slides.add_slide(layout)
    holders = slide.shapes.placeholders

    noun, meaning = holders[10], holders[11]
    noun.text_frame.text = v["word"]
    ms = v["meaning"].split(",")
    if len(ms) > 4:
      ms = ms[:4]
    meaning.text_frame.text = "\n".join(ms) 

    s_def, s = holders[12], holders[13]
    s_def.text_frame.text = "los"
    s.text_frame.text = v["word"]

    s_undef, s = holders[14], holders[15]
    s_undef.text_frame.text = "unos"
    s.text_frame.text = v["word"]

    if v["extension"] != "":
      extension = json.loads(v["extension"].replace("\'", "\""))

      pl_def, pl = holders[16], holders[17]
      pl_def.text_frame.text = "el"
      pl.text_frame.text = extension["m"]

      pl_undef, pl = holders[18], holders[19]
      pl_undef.text_frame.text = "uno"
      pl.text_frame.text = extension["m"]

    note = slide.notes_slide
    note.notes_text_frame.text = v["word"]

  def create_noun_fpl_word(self, v):
    layout = self.prs.slide_layouts.get_by_name("Noun f vocab")
    slide = self.prs.slides.add_slide(layout)
    holders = slide.shapes.placeholders

    noun, meaning = holders[10], holders[11]
    noun.text_frame.text = v["word"]
    ms = v["meaning"].split(",")
    if len(ms) > 4:
      ms = ms[:4]
    meaning.text_frame.text = "\n".join(ms) 

    s_def, s = holders[12], holders[13]
    s_def.text_frame.text = "las"
    s.text_frame.text = v["word"]

    s_undef, s = holders[14], holders[15]
    s_undef.text_frame.text = "unas"
    s.text_frame.text = v["word"]

    if v["extension"] != "":
      extension = json.loads(v["extension"].replace("\'", "\""))

      pl_def, pl = holders[16], holders[17]
      pl_def.text_frame.text = "la"
      pl.text_frame.text = extension["f"]

      pl_undef, pl = holders[18], holders[19]
      pl_undef.text_frame.text = "una"
      pl.text_frame.text = extension["f"]

    note = slide.notes_slide
    note.notes_text_frame.text = v["word"]

  def create_adj_word(self, v):
    layout = self.prs.slide_layouts.get_by_name("Adj vocab")
    slide = self.prs.slides.add_slide(layout)
    holders = slide.shapes.placeholders

    adj, meaning = holders[10], holders[11]
    adj.text_frame.text = v["word"]
    ms = v["meaning"].split(",")
    if len(ms) > 4:
      ms = ms[:4]

    meaning.text_frame.text = "\n".join(ms)

    s_m, s_f, pl_m, pl_f = holders[12], holders[13], holders[14], holders[15]
    if v["extension"] != "":
      extension = json.loads(v["extension"].replace("\'", "\""))

      s_m.text_frame.text = extension["m"]
      s_f.text_frame.text = extension["f"]
      pl_m.text_frame.text = extension["mpl"]
      pl_f.text_frame.text = extension["fpl"]

    note = slide.notes_slide
    note.notes_text_frame.text = v["word"]

  def create_default_word(self, v):
    layout = self.prs.slide_layouts.get_by_name("Default vocab")
    slide = self.prs.slides.add_slide(layout)
    holders = slide.shapes.placeholders
   
    word = holders[10]
    meaning = holders[11]
    word.text_frame.text = v["word"] 
    ms = v["meaning"].split(",")
    if len(ms) > 4:
      ms = ms[:4]

    meaning.text_frame.text = "\n".join(ms)
 
    note = slide.notes_slide
    note.notes_text_frame.text = v["word"]

    #slide = self.prs.slides.add_slide(layout)
    #holders = slide.shapes.placeholders

    #word = holders[10]
    #explaination = holders[11]
    #word.text_frame.text = v["word"] 
    #word.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 51, 0)
    #explaination.text_frame.text = v["meaning"] 
    #note = slide.notes_slide
    #note.notes_text_frame.text = " ".join([v["word"]]*2)

  def create_single_tiempo_verb_word(self, v):
    layout = self.prs.slide_layouts.get_by_name("Verb single tiempo")
    slide = self.prs.slides.add_slide(layout)
    holders = slide.shapes.placeholders
    
    variations = json.loads(v["variations"].replace("\'", "\""))

    origin, word, meaning = holders[10], holders[11], holders[12]
    origin.text_frame.text = variations["origin"]
    word.text_frame.text = v["word"]

    ms = v["meaning"].split(",")
    if len(ms) > 4:
      ms = ms[:4]

    meaning.text_frame.text = "\n".join(ms)
 

    sign = variations["formats"][0]
    tense = sign["tense"]
    person = sign["person"]

    extension = json.loads(v["extension"].replace("\'", "\""))[tense]

    holders[13].text_frame.text = extension["yo"] if extension["yo"] != "" else " "
    holders[14].text_frame.text = extension["tú"]
    holders[15].text_frame.text = extension["él/ella/Usted"]
    holders[16].text_frame.text = extension["nosotros"]
    holders[17].text_frame.text = extension["vosotros"]
    holders[18].text_frame.text = extension["ellos/ellas/Ustedes"]

    holders[19].text_frame.text = SpanishVocabMeta.tense_info[tense]
    holders[20].text_frame.text = " ".join(["人称", person, "的变位"])

    note = slide.notes_slide
    note.notes_text_frame.text = v["word"]


  def create_multi_tiempo_verb_word(self, v):
    layout = self.prs.slide_layouts.get_by_name("Verb multi tiempo")
    slide = self.prs.slides.add_slide(layout)
    holders = slide.shapes.placeholders
    
    variations = json.loads(v["variations"].replace("\'", "\""))

    origin, word, meaning = holders[10], holders[11], holders[12]
    origin.text_frame.text = variations["origin"]
    word.text_frame.text = v["word"]

    ms = v["meaning"].split(",")
    if len(ms) > 4:
      ms = ms[:4]

    meaning.text_frame.text = "\n".join(ms)
 
    formats = variations["formats"]
    holders[13].text_frame.text = "\n".join([SpanishVocabMeta.tense_info[f["tense"]] if "tense" in f.keys() else SpanishVocabMeta.tense_info[f["format"]] for f in formats])
    holders[14].text_frame.text = "\n".join([" ".join([f["person"], "的变位"]) if "person" in f.keys() else "" for f in formats])

    note = slide.notes_slide
    note.notes_text_frame.text = v["word"]

  def create_particle_verb_word(self, v):
    layout = self.prs.slide_layouts.get_by_name("Verb participle")
    slide = self.prs.slides.add_slide(layout)
    holders = slide.shapes.placeholders
    
    variations = json.loads(v["variations"].replace("\'", "\""))

    origin, word, meaning = holders[10], holders[11], holders[12]
    origin.text_frame.text = variations["origin"]
    word.text_frame.text = v["word"]

    ms = v["meaning"].split(",")
    if len(ms) > 4:
      ms = ms[:4]

    meaning.text_frame.text = "\n".join(ms)
 
    sign = variations["formats"][0]

    holders[13].text_frame.text = SpanishVocabMeta.tense_info[sign["format"]]

    note = slide.notes_slide
    note.notes_text_frame.text = v["word"]

  def create_verb_word(self, v):
    if v["variations"] == "":
      self.create_default_word(v)
    else:
      variations = json.loads(v["variations"].replace("\'", "\""))
      if "formats" in variations.keys():
        formats = variations["formats"]
        if len(formats) > 1:
          self.create_multi_tiempo_verb_word(v)
        elif "tense" in formats[0].keys():
          self.create_single_tiempo_verb_word(v)
        elif "format" in formats[0].keys():
          self.create_particle_verb_word(v)
        else:
          self.create_default_word(v)
      else:
        self.create_default_word(v)


  def create_vocab_group(self, vocabs):
    for v in vocabs:
      pos = v["dict_pos"]
      if pos == 'm.':
        self.create_noun_m_word(v) 
      elif pos == 'f.':
        self.create_noun_f_word(v)
      elif pos == 'm.pl.':
        self.create_noun_mpl_word(v) 
      elif pos == 'f.pl.':
        self.create_noun_fpl_word(v)
      elif pos == 'adj.' and v['extension'] != "":
        self.create_adj_word(v)
      elif pos == 'adj.' and v['extension'] == "":
        pass
      elif pos == 'noun.':
        pass
      elif pos in ['verb.', 'vr.', 'vi.', 'vt.', 'aux.']:
        self.create_verb_word(v)
      else:
        pass
        #self.create_default_word(v)

  def create_vocab(self):
    """
    Create 3 pos group pages, which are noun, adj, verb, refer to SpanishVocabMeta
    """
    for index, (pos, ws) in enumerate(self.displayed_word_distribution.items()):
      if index < 3 and len(ws) > 0:
        subtitle = SpanishVocabMeta.pos_info[pos][0]
        title = SpanishVocabMeta.pos_info[pos][1].upper()
        self.create_vocab_title(index+1, title, subtitle)
        self.create_vocab_group(ws)

  def create_ending(self):
    layout = self.prs.slide_layouts.get_by_name("Thanks")
    self.prs.slides.add_slide(layout)


  def save_ppt(self, destfile):
    self.prs.save(destfile)

  def convert_to_ppt(self, destfile='test.pptx'):
    self.create_home()
    self.create_statistic()
    self.create_vocab()
    self.create_ending()

    self.save_ppt(destfile)    
    
