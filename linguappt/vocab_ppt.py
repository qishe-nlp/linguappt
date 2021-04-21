from pptx import Presentation
from linguappt.lib import readCSV
import json
from abc import abstractmethod


class VocabPPT:
  """It is designed as an abstract class to be inheritated for vocabulary ppt generation of different languages,

  Note:
    **DON'T** use this class directly. To define a VocabPPT for a lanugage

      1. Define subclass inheriting ``VocabPPT``, e.g, ``ChineseVocabPPT`` 
      2. Define class variable ``_templates``, which is a :obj:`dict`, key is template genre, value is template path
      3. Define class variable ``content_keys``, which is a :obj:`list`, containing the heads in csv file
      4. Implement methods ``_create_noun(self, v)``, ``_create_verb(self, v)``, ``_create_adj(self, v)``, where ``v`` is a vocab object, corresponding to a record (one line) in csv file
    
  """

  def __init__(self, sourcefile, title="", genre="classic"):
    """Initialize ppt object, ppt title, and content from source csv file

    Args:
      sourcefile (str): csv file, whose content is written into ppt
      title (str): title written in ppt home slide
      genre (str): ppt template style 
    """

    if self.__class__.__name__ != "VocabPPT":
      self._assert_class_variables()
      self._template = self.__class__._templates[genre]
      self._prs = Presentation(self._template)
      self._title = title
      self._sourcefile = sourcefile
      self._assert_content()
      self._divide_by_pos()
    else:
      raise TypeError(self.__class__.__doc__)

  def _assert_class_variables(self):
    cls = self.__class__
    assert cls._templates != None
    assert isinstance(cls._templates, dict)
    assert cls.content_keys != None
    assert isinstance(cls.content_keys, list)

  def _divide_by_pos(self):
    """Store words in ``self.word_distribution`` as dict, whose key is PoS (part of speech) defined in ``EnglisgVocabMeta.pos_info`` and value is a list of words.
    """

    word_distribution  = {}

    for e in self.content:
      key = self.__class__._metainfo.get_pos(e["dict_pos"])
      if key in self.__class__.ALLOWED_POSES:
        if key not in word_distribution:
          word_distribution[key] = []
        word_distribution[key].append(e)
    self.word_distribution = word_distribution


  def _assert_content(self):
    """Ensure the csv file content has keys defined in ppt class
    """

    content = readCSV(self._sourcefile) 
    for e in content:
      keys = e.keys()
      assert(len(keys) == len(self.__class__.content_keys))
      for k in keys:
        assert(k in self.__class__.content_keys)
    self.content = content

  def _create_opening(self):
    """Create home slide
    """

    layout = self._prs.slide_layouts.get_by_name("Opening for chinese")
    slide = self._prs.slides.add_slide(layout)
    holders = slide.shapes.placeholders
 
    title = holders[10] 
    title.text_frame.text = self._title
    subtitle = holders[11]
    subtitle.text_frame.text = '词汇总结'

  def _create_statistics(self):
    """Create statistic slide
    """

    layout = self._prs.slide_layouts.get_by_name("Statistics")
    slide = self._prs.slides.add_slide(layout)
    holders = slide.shapes.placeholders

    for index, (pos, ws) in enumerate(self.word_distribution.items()):
      pos_name = holders[11+2*index]
      num = holders[10+2*index]
      pos_name.text_frame.text = self.__class__._metainfo.get_pos_name(pos).upper() 
      num.text_frame.text = str(len(ws))

  def _create_ending(self):
    """Create ending slide
    """

    layout = self._prs.slide_layouts.get_by_name("Thanks")
    self._prs.slides.add_slide(layout)

  def _create_vocab_title(self, title_content, subtitle_content):
    """Create vocabulary title slide for PoS

    Args:
      title_content (str): title displayed in slide
      subtitle_content (str): subtitle displayed in slide
    """

    layout = self._prs.slide_layouts.get_by_name('Title for pos')
    slide = self._prs.slides.add_slide(layout)
    holders = slide.shapes.placeholders

    title = holders[10]
    title.text_frame.text = title_content
    subtitle = holders[11]
    subtitle.text_frame.text = subtitle_content

  def _create_default_word(self, v):
    """Create default vocabulary slide, displaying word and its meaning

    Args:
      v (dict): vocabuary object, its keys should be the same as defined in subclass variable ``content_keys``
    """

    layout = self._prs.slide_layouts.get_by_name("Default")
    slide = self._prs.slides.add_slide(layout)
    holders = slide.shapes.placeholders

    pos, noun, meaning = holders[10], holders[11], holders[12]
    pos.text_frame.text = v["dict_pos"]
    noun.text_frame.text = v["word"]
    ms = v["meaning"].split(",")
    if len(ms) > 4:
      ms = ms[:4]
    meaning.text_frame.text = "\n".join(ms) 

    note = slide.notes_slide
    note.notes_text_frame.text = v["word"]

  def _create_default_word_with_examples(self, v):
    """Create default vocabulary with examples slide, displaying word, its meaning and examples 

    Args:
      v (dict): vocabuary object, its keys should be the same as defined in subclass variable ``content_keys``
    """

    layout = self._prs.slide_layouts.get_by_name("Default with examples")
    slide = self._prs.slides.add_slide(layout)
    holders = slide.shapes.placeholders

    pos, noun, meaning = holders[10], holders[11], holders[12]
    pos.text_frame.text = v["dict_pos"]
    noun.text_frame.text = v["word"]
    ms = v["meaning"].split(",")
    if len(ms) > 4:
      ms = ms[:4]
    meaning.text_frame.text = "\n".join(ms) 

    examples = json.loads(v["examples"])
    if len(examples) >= 1:
      original, translated = holders[13], holders[14]
      ex = examples[0]
      original.text_frame.text = ex["original"]
      translated.text_frame.text = ex["translated"]
    if len(examples) >= 2:
      original, translated = holders[15], holders[16]
      ex = examples[1]
      original.text_frame.text = ex["original"]
      translated.text_frame.text = ex["translated"]

    note = slide.notes_slide
    note.notes_text_frame.text = v["word"]


  def _create_vocab_group(self, pos, vocabs):
    """Create slides for vocabularies according to PoS

    Args:
      pos (str): PoS of vocabularies
      vocabs (list of dict): vocabularies with the same PoS
    """

    for v in vocabs:
      if pos == "noun":
        self._create_noun(v) 
      elif pos == "adj":
        self._create_adj(v)
      elif pos == "verb":
        self._create_verb(v)
      #else:
      #  pass


  def _create_vocab(self):
    """Create vocab group slides, which are noun, adj, verb, etc, restrained by subclass variable ALLOWED_POSES
    """

    for pos, ws in self.word_distribution.items():
      if len(ws) > 0:
        subtitle = self.__class__._metainfo.get_pos_cn_name(pos)
        title = self.__class__._metainfo.get_pos_cn_name(pos).upper()
        self._create_vocab_title(title, subtitle)
        self._create_vocab_group(pos, ws)

  def _save_ppt(self, destfile):
    """Save ppt object into file
    """

    self._prs.save(destfile)

  def convert_to_ppt(self, destfile='test.pptx'):
    """Convert csv file containing vocabulary information into pptx file

    Args:
      destfile (str): pptx file path
    """

    self._create_opening()
    self._create_statistics()
    self._create_vocab()
    self._create_ending()

    self._save_ppt(destfile)    

  @abstractmethod
  def _create_noun(self, v):
    """Create slide for vocabulary with PoS as noun

    Args:
      v (dict): vocabuary object, its keys should be the same as defined in subclass variable ``content_keys``
    """
    pass

  @abstractmethod
  def _create_verb(self, v):
    """Create slide for vocabulary with PoS as verb 

    Args:
      v (dict): vocabuary object, its keys should be the same as defined in subclass variable ``content_keys``
    """
    pass

  @abstractmethod
  def _create_adj(self, v):
    """Create slide for vocabulary with PoS as adjective 

    Args:
      v (dict): vocabuary object, its keys should be the same as defined in subclass variable ``content_keys``
    """
    pass

