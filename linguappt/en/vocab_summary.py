from linguappt.ppt import PPT
from linguappt.filereader import readCSV
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os
import math
import json

class EnglishVocabMeta:

  pos_info = {
    'noun': ('名词', 'noun', ['n.']),
    'adj': ('形容词', 'adjective', ['adj.']),
    'verb': ('动词', 'verb', ['v.', 'vt.vi.', 'vi.', 'vt.', 'aux.', 'vi.vt.']),
    'adv': ('副词', 'adverb', ['adv.']),
    'pron': ('代词', 'pronoun', ['pron.']),
    'prep': ('前置词', '', ['prep.']),
    'other': ('其他', 'others', [])
  }

  format_info = {
    'original': "原型",
    'present_participle': "现在分词",
    'past_participle': "过去分词",
    'past_tense': "过去式",
    '3': "三单",
    'comparative': "比较级",
    'superlative': "最高级",
    'singular': "单数",
    'plural': "复数"
  }

class EnglishVocabPPT(PPT):

  template_dir = os.path.dirname(__file__)
  templates = {
    "classic": os.path.join(template_dir, 'templates/vocab_english_classic.pptx'),
    "watermark": os.path.join(template_dir, 'templates/vocab_english_watermark.pptx')
  }
  lang = 'en'

  def __init__(self, sourcefile, title="", genre="classic"):
    """
    word_distribution is a list of dict
    e.g,
    word_distribution = [
      {word: 'Hello', pos: 'INTJ', num: 2},
      {word: 'good', pos: 'ADJ', num: 20},
      {word: 'then', pos: 'ADV', num: 4},
    ]
    """

    self.template = EnglishVocabPPT.templates[genre]

    self.keys = ['num', 'word', 'pos', 'meaning', 'dict_pos', 'from', 'extension', 'variations', 'examples']
    self.content = self.assert_content(sourcefile)

    self.divide_by_pos()
    self.cal_displayed_word_distribution()
    self.cal_distribution()

    self.prs = Presentation(self.template)
    self.title = title
    
  def assert_content(self, sourcefile):
    """
      Ensure the content has keys defined in ppt class
    """
    content = readCSV(sourcefile) 
    for e in content:
      keys = e.keys()
      assert(len(keys) == len(self.keys))
      for k in keys:
        assert(k in self.keys)
    return content

  def divide_by_pos(self):
    """
      Store words in self.word_distribution as dict, whose key is dict_pos and value is a list of words
      dict_pos is defined in x2cdict package
    """
    self.word_distribution  = {}

    for e in self.content:
      key = e["dict_pos"]
      if key not in self.word_distribution:
        self.word_distribution[key] = []
      self.word_distribution[key].append(e)

  def cal_displayed_word_distribution(self):
    """
      Bridge dict_pos and pos defined in EnglishVocabMeta
      Store words in self.displayed_word_distribution as dict, which are used to presented in ppt.
      Key in self.displayed_word_distribution is pos defined EnglishVocabMeta and value is a list of words
    """
    self.displayed_word_distribution = {}
    for key in EnglishVocabMeta.pos_info.keys():
      poses = EnglishVocabMeta.pos_info[key][2] 
      self.displayed_word_distribution[key] = []
      for pos in poses:
        if pos in self.word_distribution.keys():
          self.displayed_word_distribution[key].extend(self.word_distribution[pos])


  def cal_distribution(self):
    """
      Calculate the number of words according to pos defined in EnglishVocabMeta
    """
    self.distribution = [{"pos": pos, "name": EnglishVocabMeta.pos_info[pos][1], "num": len(ws)} for pos, ws in self.displayed_word_distribution.items()]

  def create_home(self):
    layout = self.prs.slide_layouts.get_by_name("Title and subtitle for chinese")
    slide = self.prs.slides.add_slide(layout)
    holders = slide.shapes.placeholders
 
    title = holders[10] 
    title.text_frame.text = self.title
    subtitle = holders[11]
    subtitle.text_frame.text = '词汇总结'

  def create_statistic(self):
    layout = self.prs.slide_layouts.get_by_name("Word count")
    slide = self.prs.slides.add_slide(layout)
    holders = slide.shapes.placeholders

    for index, info in enumerate(self.distribution[:3]):
      pos = holders[11+2*index]
      num = holders[10+2*index]
      pos.text_frame.text = info["name"].upper() 
      num.text_frame.text = str(info["num"])


  def create_vocab_title(self, pos, title_content, subtitle_content):
    layout = self.prs.slide_layouts.get_by_name('Title for '+ pos)
    slide = self.prs.slides.add_slide(layout)
    holders = slide.shapes.placeholders

    title = holders[10]
    title.text_frame.text = title_content
    subtitle = holders[11]
    subtitle.text_frame.text = subtitle_content
    note = slide.notes_slide
    note.notes_text_frame.text = title_content.lower()

  def create_rich_noun(self, v):
    layout = self.prs.slide_layouts.get_by_name("Noun with extension and examples")
    slide = self.prs.slides.add_slide(layout)
    holders = slide.shapes.placeholders

    pos, noun, meaning = holders[10], holders[11], holders[12]
    description = EnglishVocabMeta.format_info["singular"]
    if v["variations"] != "":
      variations = json.loads(v["variations"])
      description = " ".join([EnglishVocabMeta.format_info[e] for e in variations["formats"]])

    pos.text_frame.text = " ".join([v["dict_pos"], description])
    noun.text_frame.text = v["word"]
    ms = v["meaning"].split(",")
    if len(ms) > 4:
      ms = ms[:4]
    meaning.text_frame.text = "\n".join(ms) 

    extension = json.loads(v["extension"])
    single, plural = holders[13], holders[14]
    single.text_frame.text = extension["singular"]
    plural.text_frame.text = extension["plural"]

    examples = json.loads(v["examples"])
    if len(examples) == 2:
      original, translated = holders[15], holders[16]
      ex = examples[0]
      original.text_frame.text = ex["original"]
      translated.text_frame.text = ex["translated"]

      original, translated = holders[17], holders[18]
      ex = examples[1]
      original.text_frame.text = ex["original"]
      translated.text_frame.text = ex["translated"]

    note = slide.notes_slide
    note.notes_text_frame.text = v["word"]

  def create_noun_with_extension(self, v):
    layout = self.prs.slide_layouts.get_by_name("Noun with extension")
    slide = self.prs.slides.add_slide(layout)
    holders = slide.shapes.placeholders

    pos, noun, meaning = holders[10], holders[11], holders[12]
    description = EnglishVocabMeta.format_info["singular"]
    if v["variations"] != "":
      variations = json.loads(v["variations"])
      description = " ".join([EnglishVocabMeta.format_info[e] for e in variations["formats"]])
    pos.text_frame.text = " ".join([v["dict_pos"], description])

    noun.text_frame.text = v["word"]
    ms = v["meaning"].split(",")
    if len(ms) > 4:
      ms = ms[:4]
    meaning.text_frame.text = "\n".join(ms) 

    extension = json.loads(v["extension"])
    single, plural = holders[13], holders[14]
    single.text_frame.text = extension["singular"]
    plural.text_frame.text = extension["plural"]

    note = slide.notes_slide
    note.notes_text_frame.text = v["word"]

  def create_noun_with_examples(self, v):
    layout = self.prs.slide_layouts.get_by_name("Noun with examples")
    slide = self.prs.slides.add_slide(layout)
    holders = slide.shapes.placeholders

    pos, noun, meaning = holders[10], holders[11], holders[12]
    pos.text_frame.text = v["dict_pos"]
    noun.text_frame.text = v["word"]
    ms = v["meaning"].split(",")
    if len(ms) > 4:
      ms = ms[:4]
    meaning.text_frame.text = "\n".join(ms) 

    examples = json.loads(v["examples"])
    if len(examples) >= 1:
      original, translated = holders[13], holders[14]
      ex = examples[0]
      original.text_frame.text = ex["original"]
      translated.text_frame.text = ex["translated"]
    if len(examples) >= 2:
      original, translated = holders[15], holders[16]
      ex = examples[1]
      original.text_frame.text = ex["original"]
      translated.text_frame.text = ex["translated"]

    note = slide.notes_slide
    note.notes_text_frame.text = v["word"]

  def create_default_noun(self, v):
    layout = self.prs.slide_layouts.get_by_name("Default noun")
    slide = self.prs.slides.add_slide(layout)
    holders = slide.shapes.placeholders

    pos, noun, meaning = holders[10], holders[11], holders[12]
    pos.text_frame.text = "n."
    noun.text_frame.text = v["word"]
    ms = v["meaning"].split(",")
    if len(ms) > 4:
      ms = ms[:4]
    meaning.text_frame.text = "\n".join(ms) 

    note = slide.notes_slide
    note.notes_text_frame.text = v["word"]


  def create_noun_word(self, v):
    extension, variations, examples = v["extension"], v["variations"], v["examples"]
    if extension != "" and variations != "" and examples != "[]":
      self.create_rich_noun(v)
    elif extension != "" and examples == "[]":
      self.create_noun_with_extension(v)
    elif extension == "" and examples != "[]":
      self.create_noun_with_examples(v)
    else:
      self.create_default_noun(v)

  def create_rich_adj(self, v):
    layout = self.prs.slide_layouts.get_by_name("Adjective with extension and examples")
    slide = self.prs.slides.add_slide(layout)
    holders = slide.shapes.placeholders

    pos, adj, meaning = holders[10], holders[11], holders[12]
    description = ""
    if v["variations"] != "":
      variations = json.loads(v["variations"])
      description = " ".join([EnglishVocabMeta.format_info[e] for e in variations["formats"]])
    pos.text_frame.text = " ".join([v["dict_pos"], description])
    adj.text_frame.text = v["word"]
    ms = v["meaning"].split(",")
    if len(ms) > 4:
      ms = ms[:4]

    meaning.text_frame.text = "\n".join(ms)

    much, more, most = holders[12], holders[13], holders[14]
    extension = json.loads(v["extension"])

    much.text_frame.text = extension['original'] 
    more.text_frame.text = extension['comparative']
    most.text_frame.text = extension['superlative']

    examples = json.loads(v["examples"])
    if len(examples) >= 1:
      original, translated = holders[15], holders[16]
      ex = examples[0]
      original.text_frame.text = ex["original"]
      translated.text_frame.text = ex["translated"]
    if len(examples) >= 2:
      original, translated = holders[17], holders[18]
      ex = examples[1]
      original.text_frame.text = ex["original"]
      translated.text_frame.text = ex["translated"]
 
    note = slide.notes_slide
    note.notes_text_frame.text = v["word"]

  def create_adj_with_examples(self, v):
    layout = self.prs.slide_layouts.get_by_name("Adjective with examples")
    slide = self.prs.slides.add_slide(layout)
    holders = slide.shapes.placeholders

    pos, adj, meaning = holders[10], holders[11], holders[12]
    pos.text_frame.text = v["dict_pos"] 
    adj.text_frame.text = v["word"]
    ms = v["meaning"].split(",")
    if len(ms) > 4:
      ms = ms[:4]

    meaning.text_frame.text = "\n".join(ms)

    examples = json.loads(v["examples"])
    if len(examples) >= 1:
      original, translated = holders[13], holders[14]
      ex = examples[0]
      original.text_frame.text = ex["original"]
      translated.text_frame.text = ex["translated"]
    if len(examples) >= 2:
      original, translated = holders[15], holders[16]
      ex = examples[1]
      original.text_frame.text = ex["original"]
      translated.text_frame.text = ex["translated"]
 
    note = slide.notes_slide
    note.notes_text_frame.text = v["word"]



  def create_adj_with_extension(self, v):
    layout = self.prs.slide_layouts.get_by_name("Adjective with extension")
    slide = self.prs.slides.add_slide(layout)
    holders = slide.shapes.placeholders

    pos, adj, meaning = holders[10], holders[11], holders[12]
    description = ""
    if v["variations"] != "":
      variations = json.loads(v["variations"])
      description = " ".join([EnglishVocabMeta.format_info[e] for e in variations["formats"]])
    pos.text_frame.text = " ".join([v["dict_pos"], description])
    adj.text_frame.text = v["word"]
    ms = v["meaning"].split(",")
    if len(ms) > 4:
      ms = ms[:4]

    meaning.text_frame.text = "\n".join(ms)

    much, more, most = holders[13], holders[14], holders[15]
    extension = json.loads(v["extension"])

    much.text_frame.text = extension['original'] 
    more.text_frame.text = extension['comparative']
    most.text_frame.text = extension['superlative']

    note = slide.notes_slide
    note.notes_text_frame.text = v["word"]

  def create_default_adj(self, v):
    layout = self.prs.slide_layouts.get_by_name("Default adjective")
    slide = self.prs.slides.add_slide(layout)
    holders = slide.shapes.placeholders

    pos, adj, meaning = holders[10], holders[11], holders[12]
    pos.text_frame.text = v["dict_pos"] 

    adj.text_frame.text = v["word"]
    ms = v["meaning"].split(",")
    if len(ms) > 4:
      ms = ms[:4]

    meaning.text_frame.text = "\n".join(ms)

    note = slide.notes_slide
    note.notes_text_frame.text = v["word"]


  def create_adj_word(self, v):
    extension, variations, examples = v["extension"], v["variations"], v["examples"]
    if extension != "" and variations != "" and examples != "[]":
      self.create_rich_adj(v)
    elif extension != "" and examples == "[]":
      self.create_adj_with_extension(v)
    elif extension == "" and examples != "[]":
      self.create_adj_with_examples(v)
    else:
      self.create_default_adj(v)

  def create_rich_verb(self, v):
    layout = self.prs.slide_layouts.get_by_name("Verb with extension and examples")
    slide = self.prs.slides.add_slide(layout)
    holders = slide.shapes.placeholders

    pos, verb, meaning = holders[10], holders[11], holders[12]
    description = ""
    if v["variations"] != "":
      variations = json.loads(v["variations"])
      description = " ".join([EnglishVocabMeta.format_info[e] for e in variations["formats"]])
    pos.text_frame.text = " ".join([v["dict_pos"], description])
    verb.text_frame.text = v["word"]
    ms = v["meaning"].split(",")
    if len(ms) > 4:
      ms = ms[:4]

    meaning.text_frame.text = "\n".join(ms)

    original, present_participle, past_tense, past_participle, third  = holders[13], holders[14], holders[15], holders[16], holders[17]
    extension = json.loads(v["extension"])

    original.text_frame.text = extension['original'] if 'original' in extension else " "
    past_tense.text_frame.text = extension['past_tense'] if 'past_tense' in extension else " "
    third.text_frame.text = extension['3'] if '3' in extension else " "
    present_participle.text_frame.text = extension['present_participle']  if 'present_participle' in extension else " "
    past_participle.text_frame.text = extension['past_participle'] if 'past_participle' in extension else " "

    examples = json.loads(v["examples"])
    if len(examples) >= 1:
      original, translated = holders[18], holders[19]
      ex = examples[0]
      original.text_frame.text = ex["original"]
      translated.text_frame.text = ex["translated"]
    if len(examples) >= 2:
      original, translated = holders[20], holders[21]
      ex = examples[1]
      original.text_frame.text = ex["original"]
      translated.text_frame.text = ex["translated"]
 
    note = slide.notes_slide
    note.notes_text_frame.text = v["word"]

  def create_verb_with_examples(self, v):
    layout = self.prs.slide_layouts.get_by_name("Verb with examples")
    slide = self.prs.slides.add_slide(layout)
    holders = slide.shapes.placeholders

    pos, verb, meaning = holders[10], holders[11], holders[12]
    pos.text_frame.text = v['dict_pos']
    verb.text_frame.text = v["word"]
    ms = v["meaning"].split(",")
    if len(ms) > 4:
      ms = ms[:4]

    meaning.text_frame.text = "\n".join(ms)

    examples = json.loads(v["examples"])
    if len(examples) >= 1:
      original, translated = holders[13], holders[14]
      ex = examples[0]
      original.text_frame.text = ex["original"]
      translated.text_frame.text = ex["translated"]
    if len(examples) >= 2:
      original, translated = holders[15], holders[16]
      ex = examples[1]
      original.text_frame.text = ex["original"]
      translated.text_frame.text = ex["translated"]
 
    note = slide.notes_slide
    note.notes_text_frame.text = v["word"]



  def create_verb_with_extension(self, v):
    layout = self.prs.slide_layouts.get_by_name("Verb with extension")
    slide = self.prs.slides.add_slide(layout)
    holders = slide.shapes.placeholders

    pos, verb, meaning = holders[10], holders[11], holders[12]
    description = ""
    if v["variations"]:
      variations = json.loads(v["variations"])
      description = " ".join([EnglishVocabMeta.format_info[e] for e in variations["formats"]])
    pos.text_frame.text = " ".join([v["dict_pos"], description])
    verb.text_frame.text = v["word"]
    ms = v["meaning"].split(",")
    if len(ms) > 4:
      ms = ms[:4]

    meaning.text_frame.text = "\n".join(ms)

    original, past_tense, present_participle, past_participle, third = holders[13], holders[14], holders[15], holders[16], holders[17]
    extension = json.loads(v["extension"])

    original.text_frame.text = extension['original'] if 'original' in extension else " "
    past_tense.text_frame.text = extension['past_tense'] if 'past_tense' in extension else " "
    third.text_frame.text = extension['3'] if '3' in extension else " "
    present_participle.text_frame.text = extension['present_participle']  if 'present_participle' in extension else " "
    past_participle.text_frame.text = extension['past_participle'] if 'past_participle' in extension else " "

    note = slide.notes_slide
    note.notes_text_frame.text = v["word"]


  def create_default_verb(self, v):
    layout = self.prs.slide_layouts.get_by_name("Default adjective")
    slide = self.prs.slides.add_slide(layout)
    holders = slide.shapes.placeholders

    pos, verb, meaning = holders[10], holders[11], holders[12]
    pos.text_frame.text = v["dict_pos"] 

    verb.text_frame.text = v["word"]
    ms = v["meaning"].split(",")
    if len(ms) > 4:
      ms = ms[:4]

    meaning.text_frame.text = "\n".join(ms)

    note = slide.notes_slide
    note.notes_text_frame.text = v["word"]



  def create_verb_word(self, v):
    extension, variations, examples = v["extension"], v["variations"], v["examples"]
    if extension != "" and variations != "" and examples != "[]":
      self.create_rich_verb(v)
    elif extension != "" and examples == "[]":
      self.create_verb_with_extension(v)
    elif extension == "" and examples != "[]":
      self.create_verb_with_examples(v)
    else:
      self.create_default_verb(v)

  def create_vocab_group(self, vocabs):
    for v in vocabs:
      pos = v["dict_pos"]
      if pos in EnglishVocabMeta.pos_info['noun'][2]:
        self.create_noun_word(v) 
      elif pos in EnglishVocabMeta.pos_info['adj'][2]:
        self.create_adj_word(v)
      elif pos in EnglishVocabMeta.pos_info['verb'][2]:
        self.create_verb_word(v)
      else:
        pass
        #self.create_default_word(v)

  def create_vocab(self):
    """
    Create 3 pos group pages, which are noun, adj, verb, refer to EnglishVocabMeta
    """
    for index, (pos, ws) in enumerate(self.displayed_word_distribution.items()):
      if index < 3 and len(ws) > 0:
        subtitle = EnglishVocabMeta.pos_info[pos][0]
        title = EnglishVocabMeta.pos_info[pos][1].upper()
        self.create_vocab_title(pos, title, subtitle)
        self.create_vocab_group(ws)

  def create_ending(self):
    layout = self.prs.slide_layouts.get_by_name("Thanks")
    self.prs.slides.add_slide(layout)


  def save_ppt(self, destfile):
    self.prs.save(destfile)

  def convert_to_ppt(self, destfile='test.pptx'):
    self.create_home()
    self.create_statistic()
    self.create_vocab()
    self.create_ending()

    self.save_ppt(destfile)    
    
