from pptx import Presentation
from linguappt.lib import readCSV
import json
from abc import abstractmethod

class PhrasePPT:
  """It is designed as an abstract class to be inheritated for phrase ppt generation of different languages,

  Note:
    **DON'T** use this class directly. To define a PhrasePPT for a lanugage

      1. Define subclass inheriting ``PhrasePPT``, e.g, ``ChinesePhrasePPT`` 
      2. Define class variable ``_templates``, which is a :obj:`dict`, key is template genre, value is template path
      3. Define class varibale ``content_keys``, which is a :obj:`list`, containing the heads in csv file
      4. Implement method ``_create_phrase(self)``
 
  """

  def __init__(self, sourcefile, title="", genre="classic"):
    """Initialize ppt object, ppt title, and content from source csv file
    Args:
      sourcefile (str): csv file, whose content is written into ppt
      title (str): title written in ppt home slide
      genre (str): ppt template style 
    """

    if self.__class__.__name__ != "PhrasePPT":
      self._assert_class_variables()
      self._template = self.__class__._templates[genre]
      self._prs = Presentation(self._template)
      self._title = title
      self._sourcefile = sourcefile
      self._assert_content()
    else:
      raise TypeError(self.__class__.__doc__)

  def _assert_class_variables(self):
    cls = self.__class__
    assert cls._templates != None
    assert isinstance(cls._templates, dict)
    assert cls.content_keys != None
    assert isinstance(cls.content_keys, list)


  def _assert_content(self):
    """Ensure the csv file content has keys defined in ppt class
    """

    content = readCSV(self._sourcefile) 
    for e in content:
      keys = e.keys()
      assert(len(keys) == len(self.__class__.content_keys))
      for k in keys:
        assert(k in self.__class__.content_keys)
    self.content = content


  def _create_opening(self):
    """Create home slide
    """

    layout = self._prs.slide_layouts.get_by_name("Opening for chinese")

    slide = self._prs.slides.add_slide(layout)
    holders = slide.shapes.placeholders
 
    title = holders[10] 
    title.text_frame.text = self._title
    subtitle = holders[11]
    subtitle.text_frame.text = '短语总结'

  def _create_ending(self):
    """Create ending slide
    """

    layout = self._prs.slide_layouts.get_by_name("Thanks")
    self._prs.slides.add_slide(layout)

  def _save_ppt(self, destfile):
    """Save ppt object into file
    """

    self._prs.save(destfile)

  def convert_to_ppt(self, destfile='test.pptx'):
    """Convert csv file containing vocabulary information into pptx file

    Args:
      destfile (str): pptx file path
    """

    self._create_opening()
    self._create_phrase()
    self._create_ending()

    self._save_ppt(destfile)    

  @abstractmethod
  def _create_phrase(self):
    """Create phrase slide for each sentence
    """
    pass

