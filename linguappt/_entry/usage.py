from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

import os
import click

# TODO : may need to be updated
def _test_template_usage():
  template_to_be_tested = './linguappt/es/templates/vocab_spanish_classic.pptx'
  prs = Presentation(template_to_be_tested)
  templates = prs.slide_layouts
  
  # template idx:0 - name:Title and subtitle
  layout = templates.get_by_name("Title and subtitle") 
  slide = prs.slides.add_slide(layout)
  holders = slide.shapes.placeholders


  title = holders[10] 
  title.text_frame.text = 'Extra E1-1'
  subtitle = holders[11]
  subtitle.text_frame.text = '解析'

  # template idx:2 - name:Title 1
  layout = templates.get_by_name("Title 1")
  slide = prs.slides.add_slide(layout)
  holders = slide.shapes.placeholders


  title = holders[10]
  title.text_frame.text = '名词'
  subtitle = holders[11]
  subtitle.text_frame.text = 'El nombre'
  note = slide.notes_slide
  note.notes_text_frame.text = "xxxxxxxxxxx" 


  # template 2
  slide = prs.slides.add_slide(templates[2])
  holders = slide.shapes.placeholders


  word = holders[10]
  explaination = holders[11]
  word.text_frame.text = 'casa'
  explaination.text_frame.text = '家'

  # template 3
  slide = prs.slides.add_slide(templates[3])
  holders = slide.shapes.placeholders

  word = holders[10]
  explaination = holders[11]
  word.text_frame.text = 'niña'
  word.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 255)
  explaination.text_frame.text = '小女孩'


  # template 4
  slide = prs.slides.add_slide(templates[4])
  holders = slide.shapes.placeholders


  for index in range(9):
    word = holders[10+2*index]
    explaination = holders[11+2*index]
    word.text_frame.text = str(index)
    explaination.text_frame.text = str(index) 

  # template 5
  slide = prs.slides.add_slide(templates[13])
  holders = slide.shapes.placeholders

  for e in holders:
    print('%d %s' % (e.placeholder_format.idx, e.name))
 
  for index in range(4):
    word = holders[10+2*index]
    explaination = holders[11+2*index]
    word.text_frame.text = str(index*10)
    explaination.text_frame.text = str(index*10) 

  prs.save('test.pptx')

def _default():
  prs = Presentation()
  blank_slide_layout = prs.slide_layouts[6]
  slide = prs.slides.add_slide(blank_slide_layout)

  left = top = width = height = Inches(1)
  txBox = slide.shapes.add_textbox(left, top, width, height)
  tf = txBox.text_frame

  tf.text = "This is text inside a textbox"

  p = tf.add_paragraph()
  p.text = "This is a second paragraph that's bold"
  p.font.bold = True

  p = tf.add_paragraph()
  p.text = "This is a third paragraph that's big"
  p.font.size = Pt(40)

  prs.save('test.pptx')

